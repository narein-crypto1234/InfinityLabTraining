from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

# Create presentation
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "MoldVision AI: Project Demo"
slide.placeholders[1].text = "AI + Computer Vision for Mold Detection\nPresented on: " + datetime.now().strftime('%Y-%m-%d')

# Slide: Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Overview"
content = slide.placeholders[1]
content.text = (
    "This project combines Large Language Models (LLM) and Vision-Language Models (VLM) using YOLOv8\n"
    "to detect mold presence on surfaces such as bread and PCB boards.\n\n"
    "LLM Part: Predicts mold risk based on temperature, humidity, storage duration, and surface type.\n"
    "VLM Part: Uses YOLOv8 to detect mold presence in images, including real-time webcam captures.\n"
    "Results are logged and visualized for analysis."
)

# Slide: Detection Results
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Mold Detection Summary"
slide.shapes.add_picture("mold_detection_pie_chart.png", Inches(1), Inches(1.5), width=Inches(4.5))

# Slide: Confidence Graph
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Confidence Over Time"
slide.shapes.add_picture("confidence_over_time.png", Inches(1), Inches(1.5), width=Inches(6))

# Slide: How It Works
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "How It Works (YOLOv8)"
slide.placeholders[1].text = (
    "1. Image is captured from webcam or selected from storage.\n"
    "2. YOLOv8 model processes the image to detect mold presence.\n"
    "3. Bounding boxes highlight detected areas (if any).\n"
    "4. Detection results are logged to CSV and visualized."
)

# Slide: Flexibility
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Modular Detection Modes"
slide.placeholders[1].text = (
    "✅ Scan All Images\n"
    "📅 Scan Today's Captures\n"
    "📸 Scan Only Latest Image\n\n"
    "Each detection mode is accessible via its own Python script."
)

# Slide: Final Thoughts
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Final Thoughts"
slide.placeholders[1].text = (
    "This mini project showcases practical AI application in Computer Vision with flexible scanning options,\n"
    "easy reporting, and integration potential with IoT for future upgrades.\n\n"
    "Built using Python, OpenCV, Pandas, Matplotlib, and Ultralytics YOLOv8."
)

# Save the presentation
prs.save("MoldVision_AI_Presentation.pptx")
